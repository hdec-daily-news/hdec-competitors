# -*- coding: utf-8 -*-
"""슬라이드 마스터/레이아웃까지 다 검사해서 워터마크/placeholder 흔적 찾기"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

WORK = r"C:\temp\pptx_work\work.pptx"
FILLED = r"C:\temp\pptx_work\work_filled.pptx"

def dump_shapes(label, shapes):
    for i, shape in enumerate(shapes):
        try:
            print(f"  [{label}.{i}] type={shape.shape_type} name={shape.name!r} pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})")
        except Exception as e:
            print(f"  [{label}.{i}] (속성읽기 실패: {e})")
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                txt = ''.join(r.text for r in p.runs)
                if txt.strip():
                    print(f"      p{pi}: {txt[:120]!r}")

def inspect(path, label):
    print(f"\n{'='*60}\n{label}: {path}\n{'='*60}")
    prs = Presentation(path)

    # 슬라이드 마스터
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n--- Slide Master {mi} ---")
        dump_shapes(f"M{mi}", master.shapes)
        # 레이아웃들
        for li, layout in enumerate(master.slide_layouts):
            print(f"\n  -- Layout {mi}.{li}: {layout.name!r} --")
            dump_shapes(f"M{mi}L{li}", layout.shapes)

    # 슬라이드 본체
    for si, slide in enumerate(prs.slides):
        print(f"\n--- Slide {si+1} (layout: {slide.slide_layout.name!r}) ---")
        dump_shapes(f"S{si}", slide.shapes)

inspect(WORK, "원본 (work.pptx)")
inspect(FILLED, "채운 후 (work_filled.pptx)")
