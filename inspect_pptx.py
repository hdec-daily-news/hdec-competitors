# -*- coding: utf-8 -*-
"""대형사 동향_0410.pptx 구조 검사"""
import sys, os, glob
sys.stdout.reconfigure(encoding='utf-8')

target = r"C:\temp\pptx_work\jenspark2.pptx"
print(f"작업 대상: {target}\n")

from pptx import Presentation
from pptx.util import Emu

prs = Presentation(target)
print(f"슬라이드 수: {len(prs.slides)}")
print(f"슬라이드 크기: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch")
print(f"            ({prs.slide_width} x {prs.slide_height} EMU)")

for si, slide in enumerate(prs.slides):
    print(f"\n=== Slide {si+1} ===")
    for shi, shape in enumerate(slide.shapes):
        print(f"  [{shi}] type={shape.shape_type} name={shape.name!r}")
        try:
            print(f"       pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})")
        except:
            pass
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                txt = ''.join(r.text for r in p.runs)
                if txt.strip():
                    print(f"       p{pi}: {txt[:120]!r}")
        if shape.has_table:
            tbl = shape.table
            print(f"       table {len(tbl.rows)}x{len(tbl.columns)}")
            col_w = [c.width for c in tbl.columns]
            row_h = [r.height for r in tbl.rows]
            print(f"       cols(width)={col_w}")
            print(f"       rows(height)={row_h}")
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    txt = cell.text.replace('\n', ' | ')
                    print(f"         [{ri},{ci}] {txt[:120]!r}")
