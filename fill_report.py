# -*- coding: utf-8 -*-
"""
대형사 동향 한 장 보고서 — 인사이트 버전 (insight_ver.pptx 양식 기반)

구조:
 - 각사별 ● [핵심 이슈 제목] + dash 2~4개 인사이트 불릿
 - 우측 [회사명 | 기사원문] 링크 박스 → GitHub Pages 회사별 필터
 - 상단 ▶ AI 자동 수집 링크 (제목 클릭 시 원문 이동 안내)

베이스: C:\\temp\\pptx_work\\insight_ver.pptx
  (인사이트 텍스트만 제거하고 회사 레이블/구분선/이미지는 유지)
"""
import os, sys, re, zipfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────
# 경로
# ──────────────────────────────────────
WORK_DIR = r"C:\temp\pptx_work"
BASE = os.path.join(WORK_DIR, "insight_ver.pptx")
OUT = os.path.join(WORK_DIR, "report_filled.pptx")

# ──────────────────────────────────────
# 색상 / 폰트
# ──────────────────────────────────────
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(0x70, 0x70, 0x70)
DARK_BLUE = RGBColor(0x1F, 0x49, 0x7D)
LINK_BLUE = RGBColor(0x05, 0x63, 0xC1)

COMPANY_COLORS = {
    "삼성물산":  RGBColor(0x14, 0x28, 0xa0),
    "삼성E&A":   RGBColor(0x14, 0x28, 0xa0),
    "대우건설":  RGBColor(0x00, 0x3d, 0xa5),
    "GS건설":    RGBColor(0xed, 0x1c, 0x24),
    "DL이앤씨":  RGBColor(0xe9, 0x5c, 0x0c),
}

FONT = "맑은 고딕"
GH_PAGES = "https://hdec-daily-news.github.io/hdec-competitors/"

# ──────────────────────────────────────
# 양식 좌표 (insight_ver.pptx 기준)
# ──────────────────────────────────────
CONTENT_LEFT = 1857146       # ● 마커 left
BULLET_LEFT = 2048256        # dash 불릿 left
TITLE_LEFT = 2104949         # 핵심 제목 left
TITLE_WIDTH = 5500000
BULLET_WIDTH = 6400000
MARKER_WIDTH = 247802
MARKER_HEIGHT = 257861

RIGHT_LINK_LEFT = 7939735
RIGHT_LINK_WIDTH = 619963
RIGHT_LINK_HEIGHT = 457200

# 회사별 제목(●) top Y
COMPANY_TITLE_TOPS = {
    "삼성물산": 1666951,
    "삼성E&A":  3666744,
    "대우건설": 5952744,
    "GS건설":   8143646,
    "DL이앤씨": 10239451,
}
# 회사별 우측 링크 박스 top Y
COMPANY_LINK_TOPS = {
    "삼성물산": 2300630,
    "삼성E&A":  4443070,
    "대우건설": 6681521,
    "GS건설":   8820302,
    "DL이앤씨": 10962742,
}

COMPANIES = ["삼성물산", "삼성E&A", "대우건설", "GS건설", "DL이앤씨"]

# ──────────────────────────────────────
# 회사별 주간 인사이트 — 매주 업데이트
# 포맷: { title: 핵심 이슈 1줄, bullets: [dash 2~4개] }
# ──────────────────────────────────────
INSIGHTS = {
    "삼성물산": {
        "title": "압구정4구역 재건축 수주 유력 (단독 입찰)",
        "bullets": [
            "2차 현장설명회 단독 참석으로 수의계약 수순 돌입 전망",
            "업계 최고 수준의 금융 경쟁력(AA+ 신용등급)을 주요 강점으로 내세움",
            "압구정·목동·반포 등 주요 A급 수주전에 적극 참여 예정",
        ],
    },
    "삼성E&A": {
        "title": "중동 중심 실적·수주 동반 성장 및 뉴에너지 전환",
        "bullets": [
            "'26년 1분기 매출 2.3조, 영업이익 1,832억 예상 (전년비 동반 상승)",
            "중동 화공플랜트 수주 호조로 상반기 6조원 규모 신규 수주 달성 유력",
            "이란 관련 중동 사태 영향 제한적, 종전 시 중동 재건 특수 기대감 증가",
            "주택사업 의존도 탈피, 글로벌 에너지 인프라 및 신사업 다각화 추진",
        ],
    },
    "대우건설": {
        "title": "'26년 정비사업 수주액 2조원 돌파 (업계 1위)",
        "bullets": [
            "총 5개 사업장에서 2조 2,525억원 수주 달성하며 1분기 선두 기록",
            "경기 용인 기흥1구역, 서울 마포 성산 모아타운3구역 등 정비사업 연속 수주",
            "2분기 승부처인 '압구정·성수' 지역 수주전에 역량 집중 예정",
        ],
    },
    "GS건설": {
        "title": "초고층 시공 맞춤형 혁신 콘크리트 기술 개발 착수",
        "bullets": [
            "삼표산업·삼표시멘트·생고뱅코리아와 기술 협력 MOU 체결",
            "고층 운반 용이한 '저점성·고유동 콘크리트' 공동 개발 목표",
            "개발 기술을 향후 초고층 건축물 시공 현장에 적용하여 안정성 확보",
            "대치동(은마·미도) 등 강남권 주요 재건축 수주 대비 기술 경쟁력 강화",
        ],
    },
    "DL이앤씨": {
        "title": "압구정5구역 등 하이엔드 정비사업 수주 총력전",
        "bullets": [
            "압구정5구역(약 1.5조 규모) 입찰 참여, 현대건설과 '하이엔드 왕좌' 경쟁",
            "실적 및 금융 차별화 전략, 글로벌 기업 협업으로 새로운 강남 기준 제시",
            "국내 최초로 AI 기반 생성형 설계 프로그램(오바바쿠스) 적용 제안",
            "압구정을 비롯해 목동 6단지, 반포 등 대어급 수주전에 적극 참여",
        ],
    },
}

# ──────────────────────────────────────
# 콘텐츠 textbox 제거 (베이스에 있는 기존 insight/헤드라인/링크)
# ──────────────────────────────────────
def remove_content_shapes(slide):
    spTree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            left = shape.left
        except Exception:
            continue
        # 콘텐츠 영역 (● 마커 ~ 우측 링크까지 전부)
        if left >= 1800000:
            to_remove.append(shape._element)
    for el in to_remove:
        spTree.remove(el)
    print(f"  기존 콘텐츠 shape {len(to_remove)}개 제거")


def _new_tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return tb, tf


def add_company_insight(slide, company, insight, url):
    """회사별 insight 블록 — ● 마커 + 제목(링크) + dash 불릿들"""
    title_top = COMPANY_TITLE_TOPS[company]

    # ● 마커
    _, tf = _new_tb(slide, CONTENT_LEFT, title_top, MARKER_WIDTH, MARKER_HEIGHT)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "●"
    r.font.size = Pt(11)
    r.font.name = FONT
    r.font.bold = True
    r.font.color.rgb = DARK_BLUE

    # 핵심 제목 (클릭 가능 — GitHub Pages 회사 필터)
    _, tf = _new_tb(slide, TITLE_LEFT, title_top, TITLE_WIDTH, 295351)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = insight["title"]
    r.font.size = Pt(12)
    r.font.name = FONT
    r.font.bold = True
    r.font.color.rgb = BLACK
    if url:
        r.hyperlink.address = url

    # dash 불릿들
    bullet_top = title_top + 407823
    line_step = 333771  # ~0.37" between lines
    for i, b in enumerate(insight["bullets"]):
        _, tf = _new_tb(slide, BULLET_LEFT, bullet_top + i * line_step, BULLET_WIDTH, MARKER_HEIGHT)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"- {b}"
        r.font.size = Pt(10)
        r.font.name = FONT
        r.font.color.rgb = BLACK


def add_company_right_link(slide, company, url):
    """우측 [회사명 / 기사원문 ▶] 링크 박스"""
    top = COMPANY_LINK_TOPS[company]
    _, tf = _new_tb(slide, RIGHT_LINK_LEFT, top, RIGHT_LINK_WIDTH, RIGHT_LINK_HEIGHT)

    # 1줄: 회사명 (컬러)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = company
    r.font.size = Pt(9)
    r.font.name = FONT
    r.font.bold = True
    r.font.color.rgb = COMPANY_COLORS.get(company, BLACK)
    r.hyperlink.address = url

    # 2줄: 기사원문 ▶ (링크)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "기사원문 ▶"
    r2.font.size = Pt(8)
    r2.font.name = FONT
    r2.font.underline = True
    r2.font.color.rgb = LINK_BLUE
    r2.hyperlink.address = url


def add_top_guide_link(slide):
    """상단 우측에 'AI 자동수집 뉴스' 안내 + 통합 링크 2줄"""
    # 1줄: AI 자동수집 배너 링크
    _, tf = _new_tb(slide, 4000000, 280000, 5000000, 220000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "▶ AI 자동 수집 뉴스 — 제목/전체보기 클릭 시 원문 기사로 이동"
    r.font.size = Pt(10)
    r.font.name = FONT
    r.font.bold = True
    r.font.color.rgb = LINK_BLUE
    r.font.underline = True
    r.hyperlink.address = GH_PAGES
    print("  상단 AI 자동수집 안내 링크 추가")


def colorize_company_names(slide):
    """회사명 textbox의 텍스트를 코퍼레이트 컬러로"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                name = r.text.strip()
                if name in COMPANY_COLORS:
                    r.font.color.rgb = COMPANY_COLORS[name]
                    r.font.bold = True


def hide_master_shapes(pptx_path):
    """슬라이드 마스터 워터마크 숨김 (showMasterSp='0')"""
    tmp = pptx_path + ".tmp"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                    xml = data.decode("utf-8")
                    if "showMasterSp=" in xml:
                        xml = re.sub(r'showMasterSp="[01]"', 'showMasterSp="0"', xml)
                    else:
                        xml = re.sub(
                            r'<p:sld(\s+[^>]*?)>',
                            r'<p:sld\1 showMasterSp="0">',
                            xml, count=1
                        )
                    data = xml.encode("utf-8")
                zout.writestr(item, data)
    os.replace(tmp, pptx_path)
    print("  마스터 도형 숨김 완료")


def hash_name(company):
    return company.replace("&", "")


def main():
    print(f"베이스 로드: {BASE}")
    prs = Presentation(BASE)
    slide = prs.slides[0]
    print(f"  {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch, {len(slide.shapes)} shapes")

    print("\n기존 콘텐츠 shape 삭제...")
    remove_content_shapes(slide)

    print("\n인사이트 블록 + 우측 링크 추가...")
    for c in COMPANIES:
        url = f"{GH_PAGES}#{hash_name(c)}"
        add_company_insight(slide, c, INSIGHTS[c], url)
        add_company_right_link(slide, c, url)
        print(f"  [{c}] 인사이트 {len(INSIGHTS[c]['bullets'])}줄 + 링크 OK")

    print("\n상단 안내 링크...")
    add_top_guide_link(slide)

    print("\n회사명 컬러 적용...")
    colorize_company_names(slide)

    prs.save(OUT)
    print(f"\n저장: {OUT}")
    hide_master_shapes(OUT)
    print("\n완료!")


if __name__ == "__main__":
    main()
