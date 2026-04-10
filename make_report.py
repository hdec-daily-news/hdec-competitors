"""
대형사 동향 - 주간 PPT 보고서 생성
양식 (1).pptx 스타일 그대로 재현
A4 세로, 검은 헤더 블록, 파란색 강조, 회색 출처
"""

import csv
import os
from datetime import datetime, timedelta, timezone
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

KST = timezone(timedelta(hours=9))
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
CSV_PATH = os.path.join(OUTPUT_DIR, "대형사_동향.csv")
PPT_PATH = os.path.join(OUTPUT_DIR, "대형사_동향.pptx")

COMPANIES = ["삼성물산", "삼성E&A", "대우건설", "GS건설", "DL이앤씨"]

# 색상 정의 (양식과 동일)
BLUE = RGBColor(0x00, 0x00, 0xFF)
RED = RGBColor(0xFF, 0x00, 0x00)
GRAY = RGBColor(0x7E, 0x7E, 0x7E)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xFF, 0x00)
LIGHT_BLUE_BG = RGBColor(0xDC, 0xE6, 0xF1)

# 폰트 (현대하모닉 대체)
FONT_L = "맑은 고딕"
FONT_M = "맑은 고딕"


def load_articles():
    company_articles = {name: [] for name in COMPANIES}
    if not os.path.exists(CSV_PATH):
        print(f"  [경고] CSV 파일 없음: {CSV_PATH}")
        return company_articles
    with open(CSV_PATH, "r", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            company = row.get("company", "")
            if company in company_articles:
                company_articles[company].append(row)
    return company_articles


def set_shape_no_line(shape):
    """도형 테두리 제거"""
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size, font_color,
                font_name=FONT_L, bold=False, alignment=PP_ALIGN.LEFT,
                word_wrap=True, vertical_anchor=MSO_ANCHOR.TOP):
    """텍스트박스 헬퍼"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    try:
        tf.vertical_anchor = vertical_anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.name = font_name
    run.font.bold = bold
    return txBox, tf, p


def add_run(paragraph, text, font_size, font_color, font_name=FONT_L, bold=False):
    """paragraph에 run 추가"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.name = font_name
    run.font.bold = bold
    return run


def create_report(company_articles):
    now = datetime.now(KST)
    week_num = (now.day - 1) // 7 + 1
    month = now.month
    year_short = now.strftime("'%y")
    week_ago = (now - timedelta(days=7)).strftime("%m.%d")
    period = f"{week_ago}~{now.strftime('%m.%d')}"

    prs = Presentation()
    # A4 세로 (양식과 동일: 7.50 x 10.83 in)
    prs.slide_width = Emu(6858000)
    prs.slide_height = Emu(9906000)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 오른쪽 세로 파란 텍스트: "대형사 동향"
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    txBox = slide.shapes.add_textbox(
        Emu(6350000), Emu(1800000), Emu(360000), Emu(600000)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "대\n형\n사\n\n동\n향"
    run.font.size = Pt(10)
    run.font.color.rgb = BLUE
    run.font.name = FONT_L

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 상단 우측: 부서명 + 날짜
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    add_textbox(slide,
        Emu(4800000), Emu(200000), Emu(1800000), Emu(200000),
        f"'{now.strftime('%y')}.{now.month}.{week_num}주차",
        Pt(11), BLACK, FONT_L, bold=False, alignment=PP_ALIGN.RIGHT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 상단: 타이틀 영역 (테이블 헤더)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 테이블 기반 레이아웃
    table_left = Emu(500000)
    table_top = Emu(500000)
    table_width = Emu(5800000)

    # ── 메인 타이틀 테이블 ──
    tbl_shape = slide.shapes.add_table(
        rows=1, cols=2, left=table_left, top=table_top,
        width=table_width, height=Emu(350000)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Emu(1200000)
    tbl.columns[1].width = Emu(4600000)

    # 왼쪽: "대상"
    cell0 = tbl.cell(0, 0)
    cell0.text = ""
    p = cell0.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "대상 회사"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_M
    # 검은 배경
    tcPr = cell0._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '000000'})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

    # 오른쪽: "내용"
    cell1 = tbl.cell(0, 1)
    cell1.text = ""
    p = cell1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "주요 사업동향"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_M
    tcPr = cell1._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '000000'})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 회사별 내용 테이블
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    content_top = table_top + Emu(350000)
    row_height = Emu(1700000)  # 각 회사 행 높이

    content_tbl_shape = slide.shapes.add_table(
        rows=len(COMPANIES), cols=2,
        left=table_left, top=content_top,
        width=table_width, height=row_height * len(COMPANIES)
    )
    content_tbl = content_tbl_shape.table
    content_tbl.columns[0].width = Emu(1200000)
    content_tbl.columns[1].width = Emu(4600000)

    for ri, company_name in enumerate(COMPANIES):
        arts = company_articles.get(company_name, [])

        # ── 왼쪽 셀: 회사명 (검은 배경, 노란 텍스트) ──
        left_cell = content_tbl.cell(ri, 0)
        left_cell.text = ""
        left_cell.vertical_anchor = MSO_ANCHOR.TOP

        p = left_cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = company_name
        run.font.size = Pt(11.5)
        run.font.bold = True
        run.font.color.rgb = YELLOW
        run.font.name = FONT_M

        # 검은 배경
        tcPr = left_cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '000000'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

        # ── 오른쪽 셀: 기사 내용 ──
        right_cell = content_tbl.cell(ri, 1)
        right_cell.text = ""
        right_cell.vertical_anchor = MSO_ANCHOR.TOP

        # 연한 파란 배경 (양식과 동일)
        tcPr = right_cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'DCE6F1'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

        tf = right_cell.text_frame
        tf.word_wrap = True

        if not arts:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "해당 기간 주요 동향 없음"
            run.font.size = Pt(9)
            run.font.color.rgb = GRAY
            run.font.name = FONT_L
            run.font.italic = True
        else:
            for ai, art in enumerate(arts[:5]):
                if ai == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                title = art.get("title", "")
                date = art.get("date", "")

                # 번호 (파란색 볼드)
                add_run(p, f"① " if ai == 0 else f"{'②③④⑤'[ai-1]} ",
                        Pt(10), BLUE, FONT_M, bold=False)

                # 제목 (파란색)
                if len(title) > 65:
                    title = title[:62] + "..."
                add_run(p, title, Pt(10), BLUE, FONT_L)

                # 날짜 (회색, 작게)
                add_run(p, f"  ({date})", Pt(8.5), GRAY, FONT_L)

                p.space_after = Pt(4)
                p.line_spacing = Pt(16)

                # 요약이 있으면 다음 줄에 추가
                desc = art.get("description", "")
                if desc and ai < 3:  # 상위 3개만 요약 표시
                    p2 = tf.add_paragraph()
                    if len(desc) > 80:
                        desc = desc[:77] + "..."
                    add_run(p2, f"   → {desc}", Pt(8.5), BLACK, FONT_L)
                    p2.space_after = Pt(6)
                    p2.line_spacing = Pt(14)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 하단 푸터
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    add_textbox(slide,
        Emu(500000), prs.slide_height - Emu(350000),
        Emu(5800000), Emu(200000),
        f"대형사 동향 · 네이버 뉴스 API 기반 사업동향 자동 수집 · {now.strftime('%Y.%m.%d')}",
        Pt(7), GRAY, FONT_L, alignment=PP_ALIGN.CENTER)

    prs.save(PPT_PATH)
    print(f"  → PPT 저장 완료: {PPT_PATH}")


def main():
    print(f"\n{'='*50}")
    print(f"  대형사 동향 PPT 보고서 생성")
    print(f"{'='*50}\n")

    print("[1/2] CSV 데이터 로드 중...")
    company_articles = load_articles()
    for name, arts in company_articles.items():
        print(f"  [{name}] {len(arts)}건")

    print("[2/2] PPT 생성 중...")
    create_report(company_articles)
    print(f"\n  완료!")


if __name__ == "__main__":
    main()
